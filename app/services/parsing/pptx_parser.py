from pathlib import Path
from typing import Any

from pptx import Presentation

from app.services.parsing.base import ParsedBlock, ParsedDocument


class PptxParser:
    def parse(self, source_path: Path) -> ParsedDocument:
        presentation = Presentation(source_path)
        blocks: list[ParsedBlock] = []

        for index, slide in enumerate(presentation.slides, start=1):
            title = self._slide_title(slide)
            body_lines, chart_captions, image_placeholders = self._slide_content(slide)
            notes_text = self._notes_text(slide)
            lines = [f"Slide {index}"]
            if title:
                lines.append(f"Title: {title}")
            lines.extend(body_lines)
            if chart_captions:
                lines.append("Charts: " + " | ".join(chart_captions))
            if notes_text:
                lines.append("Notes: " + notes_text)
            if image_placeholders:
                lines.append(f"Image placeholders: {image_placeholders}")

            text = "\n".join(line for line in lines if line.strip()).strip()
            if not text:
                continue

            metadata = {
                "parser": "pptx",
                "slide_number": index,
                "slide_title": title,
                "shape_count": len(slide.shapes),
                "chart_captions": chart_captions,
                "image_placeholder_count": image_placeholders,
                "notes_available": bool(notes_text),
                "ocr_status": "not_run",
                "citation_label": f"{source_path.name} / slide {index}",
                "chunk_boundary": True,
            }
            title_path = [f"Slide {index}"]
            if title:
                title_path.append(title)
            blocks.append(
                ParsedBlock(
                    text=text,
                    block_type="slide",
                    page_number=index,
                    title_path=title_path,
                    metadata=metadata,
                )
            )

        return ParsedDocument(
            source_path=source_path,
            title=source_path.name,
            blocks=blocks,
            metadata={"parser": "pptx", "slide_count": len(presentation.slides)},
        )

    def _slide_title(self, slide: Any) -> str | None:
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            text = title_shape.text.strip()
            if text:
                return text
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    return text.splitlines()[0].strip()
        return None

    def _slide_content(self, slide: Any) -> tuple[list[str], list[str], int]:
        lines: list[str] = []
        chart_captions: list[str] = []
        image_placeholders = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                table_lines = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        table_lines.append(" | ".join(cells))
                if table_lines:
                    lines.append("Table:\n" + "\n".join(table_lines))
            if getattr(shape, "has_chart", False):
                caption = self._chart_caption(shape)
                if caption:
                    chart_captions.append(caption)
            if self._is_image_like(shape):
                image_placeholders += 1
        return lines, chart_captions, image_placeholders

    def _chart_caption(self, shape: Any) -> str | None:
        try:
            chart = shape.chart
            if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                text = chart.chart_title.text_frame.text.strip()
                return text or None
        except Exception:
            return None
        return None

    def _notes_text(self, slide: Any) -> str | None:
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text = text_frame.text.strip() if text_frame else ""
            return text or None
        except Exception:
            return None

    def _is_image_like(self, shape: Any) -> bool:
        shape_type = str(getattr(shape, "shape_type", "")).lower()
        return "picture" in shape_type or "placeholder" in shape_type
