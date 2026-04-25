from pathlib import Path

import pytest
import app.models  # noqa: F401
from openpyxl import Workbook
from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.db.base import Base
from app.models.chunk import Chunk
from app.schemas.documents import DocumentIngestRequest
from app.schemas.retrieval import RetrievalFilter, SearchRequest
from app.services.chunking.service import ChunkingService
from app.services.ingestion.service import DocumentIngestionService
from app.services.parsing.pptx_parser import PptxParser
from app.services.parsing.registry import ParserRegistry
from app.services.parsing.xlsx_parser import XlsxParser
from app.services.retrieval.service import RetrievalService
from app.services.storage.service import StoredFile


class FakeOpenSearchChunkIndexer:
    def index_chunk(self, chunk, document, version) -> bool:
        return True


@pytest.fixture()
def db_session():
    engine = create_engine("sqlite:///:memory:")
    Base.metadata.create_all(engine)
    session = sessionmaker(bind=engine)()
    try:
        yield session
    finally:
        session.close()


def test_registry_supports_xlsx_and_pptx(tmp_path: Path):
    registry = ParserRegistry()
    assert isinstance(registry.get_parser(tmp_path / "sample.xlsx"), XlsxParser)
    assert isinstance(registry.get_parser(tmp_path / "sample.pptx"), PptxParser)


def test_xlsx_parser_and_ingestion_preserve_sheet_cell_metadata(
    tmp_path: Path,
    db_session,
    monkeypatch,
):
    xlsx_path = tmp_path / "报价汇总表.xlsx"
    _write_sample_xlsx(xlsx_path)
    monkeypatch.setattr(
        "app.services.ingestion.service.OpenSearchChunkIndexer",
        lambda: FakeOpenSearchChunkIndexer(),
    )

    job = DocumentIngestionService(db_session).ingest_uploaded_file(
        DocumentIngestRequest(
            source_uri=xlsx_path.name,
            title="报价汇总表",
            source_type="enterprise",
            document_type="excel",
        ),
        StoredFile(
            storage_uri=str(xlsx_path),
            local_path=xlsx_path,
            file_name=xlsx_path.name,
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
    )

    assert job.status == "completed"
    chunks = db_session.query(Chunk).filter(Chunk.document_id == job.document_id).all()
    assert chunks
    first_metadata = chunks[0].metadata_json
    assert first_metadata["parser"] == "xlsx"
    assert first_metadata["sheet_name"] == "报价汇总"
    assert first_metadata["cell_range"].startswith("A1:")
    assert first_metadata["row_count"] >= 2
    assert first_metadata["column_count"] >= 3
    assert first_metadata["citation_label"].endswith(f"报价汇总 / {first_metadata['cell_range']}")

    response = _fallback_search(
        db_session,
        query="投标总价 付款比例",
        document_id=job.document_id,
        document_type="excel",
        monkeypatch=monkeypatch,
    )
    assert response.results
    assert response.results[0].metadata["sheet_name"] == "报价汇总"
    assert "cell_range" in response.results[0].metadata


def test_pptx_parser_and_ingestion_preserve_slide_metadata(
    tmp_path: Path,
    db_session,
    monkeypatch,
):
    pptx_path = tmp_path / "项目汇报.pptx"
    _write_sample_pptx(pptx_path)
    monkeypatch.setattr(
        "app.services.ingestion.service.OpenSearchChunkIndexer",
        lambda: FakeOpenSearchChunkIndexer(),
    )

    job = DocumentIngestionService(db_session).ingest_uploaded_file(
        DocumentIngestRequest(
            source_uri=pptx_path.name,
            title="项目汇报",
            source_type="enterprise",
            document_type="pptx",
        ),
        StoredFile(
            storage_uri=str(pptx_path),
            local_path=pptx_path,
            file_name=pptx_path.name,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    )

    assert job.status == "completed"
    chunks = db_session.query(Chunk).filter(Chunk.document_id == job.document_id).all()
    assert chunks
    first_metadata = chunks[0].metadata_json
    assert first_metadata["parser"] == "pptx"
    assert first_metadata["slide_number"] == 1
    assert first_metadata["slide_title"] == "总体建设目标"
    assert first_metadata["citation_label"].endswith("项目汇报.pptx / slide 1")

    response = _fallback_search(
        db_session,
        query="BIM 平台能力",
        document_id=job.document_id,
        document_type="pptx",
        monkeypatch=monkeypatch,
    )
    assert response.results
    assert response.results[0].metadata["slide_number"] == 1
    assert response.results[0].section_path[0] == "Slide 1"


def test_chunker_preserves_structured_block_metadata(tmp_path: Path):
    xlsx_path = tmp_path / "进度计划.xlsx"
    _write_sample_xlsx(xlsx_path)
    parsed = XlsxParser(rows_per_block=5).parse(xlsx_path)
    chunks = ChunkingService().chunk(parsed)

    assert chunks
    assert chunks[0].metadata["sheet_name"] == "报价汇总"
    assert "cell_range" in chunks[0].metadata


def _fallback_search(db_session, query: str, document_id: str, document_type: str, monkeypatch):
    def fail_sparse(*_args, **_kwargs):
        raise RuntimeError("force database fallback")

    service = RetrievalService(db_session)
    monkeypatch.setattr(service, "_execute_sparse_search", fail_sparse)
    return service.search(
        SearchRequest(
            query=query,
            retrieval_mode="sparse",
            enable_dense=False,
            filters=RetrievalFilter(document_id=document_id, document_type=document_type),
            top_k=3,
        )
    )


def _write_sample_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "报价汇总"
    sheet.append(["项目", "金额（万元）", "付款比例"])
    sheet.append(["投标总价", 1280.5, "30%"])
    sheet.append(["安装工程费", 320, "=B2*0.25"])
    workbook.save(path)


def _write_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "总体建设目标"
    slide.placeholders[1].text = "本页介绍 BIM 平台能力、实施计划和交付风险。"
    presentation.save(path)
